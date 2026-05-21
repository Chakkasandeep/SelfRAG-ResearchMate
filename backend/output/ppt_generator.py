from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import os

def generate_pptx(review: str, citations: list[dict]) -> str:
    """
    Generates a high-quality PowerPoint deck using a bespoke corporate-academic theme:
    - Slate background (0x0F, 0x17, 0x2A)
    - Emerald accent highlights (0x10, 0xB9, 0x81)
    - Pristine typography sizing and custom card panels
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)  # standard widescreen 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Use standard blank slide layout (index 6 in default templates)
    blank_layout = prs.slide_layouts[6]

    # ── TITLE SLIDE ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)  # Slate-900

    # Title header text box
    tf = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.8))
    tf_frame = tf.text_frame
    tf_frame.word_wrap = True
    p = tf_frame.paragraphs[0]
    p.text = "Literature Review"
    p.font.size = Pt(46)
    p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)  # Emerald-500
    p.font.bold = True
    p.font.name = 'Arial'

    # Subtitle info box
    sub = slide.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(11.0), Inches(1.2))
    sub_frame = sub.text_frame
    sub_frame.word_wrap = True
    p_sub = sub_frame.paragraphs[0]
    p_sub.text = f"Synthesized by ResearchMate · {len(citations)} Sources Verified\nSelf-RAG Hallucination Free Pipeline"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)  # Slate-400
    p_sub.font.name = 'Arial'

    # ── CONTENT SLIDES ───────────────────────────────────────
    # Split review text into paragraph sections
    paragraphs = review.split("\n\n")
    slide_count = 0
    
    for para in paragraphs:
        cleaned_para = para.strip()
        # Skip headers or extremely short paragraphs
        if len(cleaned_para) < 40 or cleaned_para.startswith("#"):
            continue
            
        # Limit content slides to prevent oversized files
        if slide_count >= 8:
            break
            
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

        # Slide header textbox
        header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.8))
        h_frame = header_box.text_frame
        h_p = h_frame.paragraphs[0]
        h_p.text = f"Research Synthesis — Key Insight {slide_count + 1}"
        h_p.font.size = Pt(24)
        h_p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
        h_p.font.bold = True
        h_p.font.name = 'Arial'

        # Main slide body textbox
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.0))
        c_frame = content_box.text_frame
        c_frame.word_wrap = True
        
        c_p = c_frame.paragraphs[0]
        c_p.text = cleaned_para
        c_p.font.size = Pt(15)
        c_p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)  # Slate-200
        c_p.line_spacing = 1.3
        c_p.font.name = 'Arial'
        
        slide_count += 1

    # ── REFERENCES SLIDE ─────────────────────────────────────
    if citations:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

        # References header box
        ref_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.8))
        ref_frame = ref_box.text_frame
        ref_p = ref_frame.paragraphs[0]
        ref_p.text = "Academic References"
        ref_p.font.size = Pt(24)
        ref_p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
        ref_p.font.bold = True
        ref_p.font.name = 'Arial'

        # References bullet list textbox
        list_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.2))
        list_frame = list_box.text_frame
        list_frame.word_wrap = True
        
        # Display up to 10 references to ensure perfect formatting
        for idx, c in enumerate(citations[:10]):
            p_item = list_frame.add_paragraph() if idx > 0 else list_frame.paragraphs[0]
            p_item.text = f"[{c['id']}] {c['authors']} ({c['year'] or 'n.d.'}). {c['title']}."
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)  # Slate-400
            p_item.space_after = Pt(8)
            p_item.font.name = 'Arial'

    # Save to a secure temporary path
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp_path = tmp.name
    tmp.close()  # close so pptx can write to it
    
    prs.save(tmp_path)
    return tmp_path
